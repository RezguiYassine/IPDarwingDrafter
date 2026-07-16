#!/usr/bin/env python3
"""Build the AP3 vectorization-pipeline status PowerPoint (problems, fixes,
headline metrics, hatching deep-dive, next steps). Source content is the
README's bug-fix log / project status / roadmap sections, summarized.

    python -m tools.build_status_deck --out output/presentations/AP3_status.pptx
"""
from __future__ import annotations

import argparse

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

NAVY = RGBColor(0x10, 0x2A, 0x43)
TEAL = RGBColor(0x1F, 0x7A, 0x8C)
LIGHT = RGBColor(0xF2, 0xF5, 0xF7)
DARK_TEXT = RGBColor(0x20, 0x20, 0x20)
GREEN = RGBColor(0x2E, 0x8B, 0x57)
RED = RGBColor(0xB0, 0x3A, 0x2E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _bg(slide, color):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = color


def _textbox(slide, l, t, w, h, text, size=18, bold=False, color=DARK_TEXT,
             align=PP_ALIGN.LEFT, font="Calibri", anchor=None):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    if anchor:
        tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = font
    return box


def _bullets(slide, l, t, w, h, items, size=16, color=DARK_TEXT, line_spacing=1.15):
    """items: list of (text, level, bold_prefix_or_None)"""
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        text, level = item[0], item[1]
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = level
        p.line_spacing = line_spacing
        p.space_after = Pt(6)
        bullet = "▸ " if level == 0 else "–  "
        r = p.add_run()
        r.text = bullet + text
        r.font.size = Pt(size - level * 2)
        r.font.color.rgb = color
        r.font.name = "Calibri"
    return box


def title_slide(prs, title, subtitle, footer):
    s = _blank(prs)
    _bg(s, NAVY)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(4.55), SLIDE_W, Inches(0.06))
    bar.fill.solid(); bar.fill.fore_color.rgb = TEAL; bar.line.fill.background()
    _textbox(s, Inches(0.8), Inches(2.6), Inches(11.7), Inches(1.6), title,
              size=40, bold=True, color=WHITE)
    _textbox(s, Inches(0.8), Inches(3.75), Inches(11.7), Inches(0.7), subtitle,
              size=20, color=RGBColor(0xC9, 0xD8, 0xE2))
    _textbox(s, Inches(0.8), Inches(6.7), Inches(11.7), Inches(0.5), footer,
              size=13, color=RGBColor(0x8F, 0xA6, 0xB4))
    return s


def section_slide(prs, number, title, kicker):
    s = _blank(prs)
    _bg(s, NAVY)
    _textbox(s, Inches(0.8), Inches(2.6), Inches(2.5), Inches(1.2), number,
              size=64, bold=True, color=TEAL)
    _textbox(s, Inches(0.8), Inches(3.7), Inches(11.5), Inches(1.0), title,
              size=34, bold=True, color=WHITE)
    _textbox(s, Inches(0.8), Inches(4.5), Inches(11.5), Inches(0.6), kicker,
              size=16, color=RGBColor(0xC9, 0xD8, 0xE2))
    return s


def content_header(slide, title, kicker=None):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(1.05))
    bar.fill.solid(); bar.fill.fore_color.rgb = NAVY; bar.line.fill.background()
    _textbox(slide, Inches(0.55), Inches(0.12), Inches(12.2), Inches(0.6), title,
              size=26, bold=True, color=WHITE)
    if kicker:
        _textbox(slide, Inches(0.55), Inches(0.62), Inches(12.2), Inches(0.4), kicker,
                  size=13, color=RGBColor(0xC9, 0xD8, 0xE2))


def bullet_slide(prs, title, kicker, items, size=17):
    s = _blank(prs)
    _bg(s, WHITE)
    content_header(s, title, kicker)
    _bullets(s, Inches(0.7), Inches(1.35), Inches(12.0), Inches(5.8), items, size=size)
    return s


def problem_fix_slide(prs, title, kicker, problem, fix, result, result_good=True):
    s = _blank(prs)
    _bg(s, WHITE)
    content_header(s, title, kicker)
    col_w = Inches(5.9)
    y0 = Inches(1.35)
    # Problem card
    pcard = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.55), y0, col_w, Inches(2.55))
    pcard.fill.solid(); pcard.fill.fore_color.rgb = LIGHT; pcard.line.color.rgb = RED
    pcard.line.width = Pt(1.25)
    _textbox(s, Inches(0.8), y0 + Inches(0.12), col_w - Inches(0.5), Inches(0.4),
              "PROBLEM", size=14, bold=True, color=RED)
    _bullets(s, Inches(0.8), y0 + Inches(0.55), col_w - Inches(0.5), Inches(1.9),
              [(problem, 0)], size=15)
    # Fix card
    fcard = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.85), y0, col_w, Inches(2.55))
    fcard.fill.solid(); fcard.fill.fore_color.rgb = LIGHT; fcard.line.color.rgb = TEAL
    fcard.line.width = Pt(1.25)
    _textbox(s, Inches(7.1), y0 + Inches(0.12), col_w - Inches(0.5), Inches(0.4),
              "FIX", size=14, bold=True, color=TEAL)
    _bullets(s, Inches(7.1), y0 + Inches(0.55), col_w - Inches(0.5), Inches(1.9),
              [(fix, 0)], size=15)
    # Result strip
    y1 = y0 + Inches(2.75)
    rcard = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.55), y1, Inches(12.2), Inches(2.0))
    rcard.fill.solid()
    rcard.fill.fore_color.rgb = RGBColor(0xEA, 0xF6, 0xEE) if result_good else RGBColor(0xFB, 0xEC, 0xEA)
    rcard.line.color.rgb = GREEN if result_good else RED
    rcard.line.width = Pt(1.25)
    _textbox(s, Inches(0.8), y1 + Inches(0.1), Inches(11.6), Inches(0.4),
              "RESULT", size=14, bold=True, color=GREEN if result_good else RED)
    _bullets(s, Inches(0.8), y1 + Inches(0.5), Inches(11.6), Inches(1.4),
              [(r, 0) for r in result], size=15)
    return s


def table_slide(prs, title, kicker, headers, rows, col_widths=None, font_size=13.5,
                 highlight_col=None):
    s = _blank(prs)
    _bg(s, WHITE)
    content_header(s, title, kicker)
    n_rows, n_cols = len(rows) + 1, len(headers)
    left, top = Inches(0.55), Inches(1.35)
    width, height = Inches(12.25), Inches(min(5.7, 0.5 * n_rows + 0.3))
    gtable = s.shapes.add_table(n_rows, n_cols, left, top, width, height).table
    if col_widths:
        for i, w in enumerate(col_widths):
            gtable.columns[i].width = Inches(w)
    for j, h in enumerate(headers):
        cell = gtable.cell(0, j)
        cell.text = h
        cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
        p = cell.text_frame.paragraphs[0]
        p.runs[0].font.color.rgb = WHITE
        p.runs[0].font.bold = True
        p.runs[0].font.size = Pt(font_size)
        p.alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            cell = gtable.cell(i, j)
            cell.text = str(val)
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT if i % 2 == 0 else WHITE
            p = cell.text_frame.paragraphs[0]
            p.runs[0].font.size = Pt(font_size)
            p.runs[0].font.color.rgb = DARK_TEXT
            p.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            if highlight_col is not None and j == highlight_col:
                p.runs[0].font.bold = True
                p.runs[0].font.color.rgb = GREEN
    return s


def image_pair_slide(prs, title, kicker, left_img, left_caption, right_img, right_caption,
                      source_note=None):
    s = _blank(prs)
    _bg(s, WHITE)
    content_header(s, title, kicker)
    col_w = Inches(5.95)
    top = Inches(1.45)
    max_h = Inches(4.7)

    def _place(img_path, left, caption):
        from PIL import Image as _PILImage
        with _PILImage.open(img_path) as im:
            iw, ih = im.size
        ratio = ih / iw
        w = col_w
        h = Emu(int(w * ratio))
        if h > max_h:
            h = max_h
            w = Emu(int(h / ratio))
        x = left + (col_w - w) / 2
        s.shapes.add_picture(img_path, x, top, width=w, height=h)
        cap = s.shapes.add_textbox(left, top + max_h + Inches(0.12), col_w, Inches(0.6))
        tf = cap.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = caption
        r.font.size = Pt(14); r.font.color.rgb = DARK_TEXT; r.font.bold = True

    _place(left_img, Inches(0.55), left_caption)
    _place(right_img, Inches(6.85), right_caption)
    if source_note:
        _textbox(s, Inches(0.55), Inches(6.95), Inches(12.2), Inches(0.4), source_note,
                  size=11, color=RGBColor(0x70, 0x78, 0x80))
    return s


def single_image_slide(prs, title, kicker, img_path, caption=None, source_note=None):
    s = _blank(prs)
    _bg(s, WHITE)
    content_header(s, title, kicker)
    from PIL import Image as _PILImage
    with _PILImage.open(img_path) as im:
        iw, ih = im.size
    ratio = ih / iw
    top = Inches(1.35)
    max_w, max_h = Inches(12.2), Inches(5.1)
    w = max_w
    h = Emu(int(w * ratio))
    if h > max_h:
        h = max_h
        w = Emu(int(h / ratio))
    x = (SLIDE_W - w) / 2
    s.shapes.add_picture(img_path, x, top, width=w, height=h)
    y2 = top + h + Inches(0.15)
    if caption:
        _textbox(s, Inches(0.55), y2, Inches(12.2), Inches(0.5), caption, size=14,
                  bold=True, color=DARK_TEXT, align=PP_ALIGN.CENTER)
        y2 += Inches(0.45)
    if source_note:
        _textbox(s, Inches(0.55), y2, Inches(12.2), Inches(0.35), source_note, size=11,
                  color=RGBColor(0x70, 0x78, 0x80), align=PP_ALIGN.CENTER)
    return s


def pipeline_diagram_slide(prs):
    s = _blank(prs)
    _bg(s, WHITE)
    content_header(s, "The Pipeline", "Stage 0 + four stages — raster sketch → editable SVG / DXF")
    stages = [
        ("Stage 0\nReferences", "Remove numerals,\nleaders, captions"),
        ("Stage 1\nPreprocessing", "Clean + skeletonize\n(SketchCleanNet)"),
        ("Stage 2\nStroke Extraction", "Skeleton → stroke\ngraph (+ hachures)"),
        ("Stage 3\nPrimitive Fitting", "RANSAC: line / circle /\narc / ellipse / path"),
        ("Stage 4\nExport", "SVG + DXF\n(ISO 128 layers)"),
    ]
    n = len(stages)
    box_w, box_h = Inches(2.05), Inches(1.7)
    gap = Inches(0.32)
    total_w = box_w * n + gap * (n - 1)
    x = (SLIDE_W - total_w) / 2
    y = Inches(2.6)
    centers = []
    for i, (head, sub) in enumerate(stages):
        shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, box_w, box_h)
        shp.fill.solid(); shp.fill.fore_color.rgb = NAVY if i != 4 else TEAL
        shp.line.fill.background()
        tf = shp.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p0 = tf.paragraphs[0]
        p0.alignment = PP_ALIGN.CENTER
        r0 = p0.add_run(); r0.text = head; r0.font.bold = True; r0.font.size = Pt(14)
        r0.font.color.rgb = WHITE
        p1 = tf.add_paragraph(); p1.alignment = PP_ALIGN.CENTER
        r1 = p1.add_run(); r1.text = sub; r1.font.size = Pt(11)
        r1.font.color.rgb = RGBColor(0xD8, 0xE4, 0xEA)
        centers.append((x + box_w, y + box_h / 2))
        x += box_w + gap
    for i in range(n - 1):
        ax0, ay = centers[i]
        arrow = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, ax0, ay - Inches(0.12),
                                    gap, Inches(0.24))
        arrow.fill.solid(); arrow.fill.fore_color.rgb = RGBColor(0xB0, 0xB8, 0xC0)
        arrow.line.fill.background()
    _bullets(s, Inches(0.7), Inches(4.85), Inches(12.0), Inches(1.8), [
        ("Stage 0 also reinjects references at export; Stage 2 side-layers hachures; "
         "Stage 3 reinjects them as styled primitives excluded from confidence gates.", 0),
        ("59 commits since the initial RANSAC-based pipeline — every stage has had at least one "
         "correctness or performance fix (next slides).", 0),
    ], size=15)
    return s


def closing_slide(prs):
    s = _blank(prs)
    _bg(s, NAVY)
    _textbox(s, Inches(0.8), Inches(2.9), Inches(11.7), Inches(1.0),
              "Questions / discussion", size=36, bold=True, color=WHITE)
    _textbox(s, Inches(0.8), Inches(3.9), Inches(11.7), Inches(1.2),
              "Repo: IPDarwingDrafter  ·  README.md has the full bug-fix log, metrics, and roadmap",
              size=16, color=RGBColor(0xC9, 0xD8, 0xE2))
    return s


def build(out_path, figs):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    title_slide(
        prs,
        "AP3 Vectorization Pipeline",
        "Hand-drawn / patent engineering sketch → editable SVG & DXF — status, fixes, and next steps",
        "IP DrawingDrafter · HAW Landshut · 2026-06-29",
    )

    bullet_slide(prs, "Agenda", None, [
        ("The pipeline at a glance (Stage 0 → 4)", 0),
        ("Problems found and fixed, by category", 0),
        ("Corners  ·  short-stroke fragmentation  ·  long-stroke splitting  ·  curves  ·  "
         "circles & polygons  ·  hatching", 1),
        ("Headline accuracy results (Drawing2CAD benchmark)", 0),
        ("Deep dive: hatching — the hardest remaining problem, work in progress", 0),
        ("Known limitations and next steps", 0),
    ], size=19)

    pipeline_diagram_slide(prs)

    # ---- Problems & fixes, one per category --------------------------------
    section_slide(prs, "01", "Problems & Fixes", "What broke, why, and how it was fixed — by category")

    problem_fix_slide(
        prs, "Corners", "Stage 2 — keypoint detection",
        problem="The classical CN (crossing-number) skeleton tracer detects endpoints and "
                "junctions reliably, but is noisy at sharp corners — a key driver of jagged "
                "or over-split output on cornered shapes.",
        fix="Trained a lightweight keypoint CNN (stacked hourglass, 3 classes: endpoint / "
            "junction / corner) on Drawing2CAD. Topology extraction now consumes keypoint "
            "clusters directly. Production uses a fusion seeding: CN's endpoints/junctions "
            "+ the CNN's corners.",
        result=[
            "Val peak-F1 0.83 overall, 0.96 on corners specifically.",
            "Fusion beats pure-CN on every D2C accuracy metric: Chamfer mean −1.1%, p95 −33%, "
            "IoU +2%, primitives −10%, and 6–8× faster on cornered closed shapes.",
            "Open gap: neutral (not yet a win) on real patent drawings — out-of-distribution; "
            "see Next Steps.",
        ],
    )

    problem_fix_slide(
        prs, "Short strokes — fragmentation", "Stage 2 — graph topology",
        problem="The CN map emits a phantom junction at every Zhang-Suen staircase bend and "
                "every stroke crossing, so one logical stroke gets split into many short "
                "edges/primitives — worst on dense patent scans (thousands of 2–3px stubs).",
        fix="_simplify_graph pass run to a fixed point: prune short dead-end spurs, dissolve "
            "degree-2 phantom junctions (real corners are never CN≥3, so these are always "
            "artefacts), merge collinear edges straight through real junctions.",
        result=[
            "Drawing2CAD primitives: mean 4.35 → 2.73 (−37%), p95 12 → 7.",
            "PatentData primitives: mean 2 364 → 771 (−67%); a 30-part assembly: 2 497 → 848 edges.",
            "Accuracy unchanged-to-slightly-better (IoU +0.3%, Chamfer −1.3%) — pure cleanup, no geometry moved.",
        ],
    )

    image_pair_slide(
        prs, "Short strokes — actual before/after", "Same sample, _simplify_graph off vs. on",
        str(figs / "strokes_before.png"), "Before — 936 fragments (each colour = one edge)",
        str(figs / "strokes_after.png"), "After — 311 edges, long strokes consolidated",
        source_note="Generated live from data/samples/Picture1.png through Stage 1 + Stage 2 "
                     "with only stage2.simplify_graph toggled; same skeleton both times.",
    )

    problem_fix_slide(
        prs, "Long strokes — splitting", "Stage 0 / Stage 2",
        problem="Patent reference numerals, leader lines, and figure captions enter Stage 1 "
                "as normal ink; every contact point becomes a skeleton junction, breaking up "
                "long real outlines.",
        fix="Stage 0 detects and removes leadered labels, segment-adjacent labels, compact "
            "unleadered labels, and FIG/panel captions before Stage 1 — with bounded "
            "iterative removal, leader-tip trimming, and local mask repair — then reinjects "
            "them at export (SVG crop overlays + leaders, DXF leader geometry).",
        result=[
            "Filtered 100-sample probe: mean Stage 2 edges 307.95 → 199.31 after Stage 0.",
            "Reference detector moved from a naive heuristic to OCR-based detection, fixing a "
            "catastrophic over-removal bug (65% → 7.7% over-removal of real ink).",
            "Final residual scan: 100/100 active removals, 99/100 with zero residual detections.",
        ],
    )

    problem_fix_slide(
        prs, "Curves", "Stage 2 / Stage 3 — fitting",
        problem="Long near-straight skeleton edges produced wiggly B-spline output (Runge "
                "oscillation); compound curves (fillets, isometric silhouettes, S-curves) "
                "didn't fit any single arc/ellipse and fell back to jagged raw polylines "
                "(21% of patent primitives).",
        fix="Bounding-box + on-skeleton deviation guards reject bad splines and fall back to "
            "RDP points. New compound-path fitter: corner-split the dense chain (genuine "
            "corners split, smooth curvature doesn't), fit each piece as line/arc/cubic-Bézier, "
            "emit one connected path primitive.",
        result=[
            "Patent jagged polylines: 16% → 0% of curved edges; 83% now smooth.",
            "D2C all-views Chamfer mean −39% (1.448 → 0.884 px), p95 −69% (4.75 → 1.47 px) "
            "from the spline-oscillation fix alone.",
            "Isometric-view tail (the worst outlier) collapses: worst-case Chamfer 10.3 → 1.5 px.",
        ],
    )

    problem_fix_slide(
        prs, "Circles & polygons", "Stage 3 — RANSAC fitting",
        problem="Scan noise produced small closed loops mis-fit as spurious circles; a fixed "
                "pixel-count noise filter then over-corrected and deleted genuine small "
                "circles (radius 6–13 px); right-angle rectangles had no circle/ellipse fit "
                "and fell back to a 300-point round-looking polyline.",
        fix="Circularity guard (radial-RMS deviation, size-adaptive threshold) separates real "
            "small circles from noise loops instead of a blunt size cutoff. New closed-polygon "
            "fitter tries an RDP N-vertex polygon before the polyline fallback.",
        result=[
            "Genuine small circles preserved; noise loops still removed.",
            "Rectangles now export as a closed 4-vertex polygon, not a round polyline.",
            "Fixed a related D2C zero-output bug: tiny circles (<40px) were dropped before "
            "the circularity guard ever got to evaluate them.",
        ],
    )

    problem_fix_slide(
        prs, "Hatching / hachures", "Stage 2 / Stage 3 / Stage 4",
        problem="Section hatching (~45% of all primitives in patent drawings) creates many "
                "short, parallel edges and extra junctions where hatches meet outlines — "
                "Stage 3 then faithfully (and badly) fits an already-fragmented graph.",
        fix="Diagnosed as a Stage-2 topology problem, not a fitting problem. Side-layer "
            "design: Stage 2 extracts dense local clusters of short similarly-angled strokes "
            "into removed_hachures, re-simplifies the main graph without them, and Stage 3/4 "
            "reinject them as a styled HACHURE layer excluded from main-geometry gates.",
        result=[
            "On hachure-cleaned sketches: main edges 186.2 → 64.1 mean, micro-edge ratio 26.8% → 0.0%.",
            "Median main edge length 17.2 px → 75.1 px.",
            "Parametric HATCH-region rendering exists but is currently disabled (precision risk) "
            "— see the dedicated deep-dive.",
        ],
    )

    single_image_slide(
        prs, "Hatching — actual before/after", "Real patent sample, EP3339949B1 F0002",
        str(figs / "hachure_before_after.png"),
        caption="norefs (input) → old graph (141 edges, hatching fragments outlines) → "
                "removed hachure (113 edges extracted) → new graph (25 main edges, clean outlines)",
        source_note="From output/hachure_stage2_probe_v2/ (Stage-2 side-layer extraction probe).",
    )

    # ---- Headline metrics ----------------------------------------------------
    table_slide(
        prs, "Headline accuracy — Drawing2CAD benchmark", "1 000 test-set samples, Front view, seed 42",
        ["Metric", "Before", "After all fixes", "Change"],
        [
            ["Pixel IoU — mean", "0.619", "0.681", "+10.0%"],
            ["Pixel IoU — worst 5% (p05)", "0.427", "0.553", "+30%"],
            ["Chamfer distance — mean", "3.83 px", "0.95 px", "−75%"],
            ["Chamfer distance — p95", "21.3 px", "1.84 px", "−91%"],
            ["Zero-output samples", "12 / 1000", "0 / 1000", "−100%"],
            ["Stage 2 time — mean", "13.2 s", "10.2 s", "−23%"],
            ["Primitives per sketch — mean", "4.35", "2.73", "−37%"],
        ],
        col_widths=[5.0, 2.3, 2.95, 2.0], font_size=15, highlight_col=3,
    )

    table_slide(
        prs, "Where the remaining error comes from", "Per-stage Chamfer attribution, 120 D2C test samples",
        ["Stage", "Before B-spline fix", "After B-spline fix"],
        [
            ["Stage 1 (preprocessing)", "0.000 px", "0.000 px"],
            ["Stage 2 (stroke extraction)", "0.003 px", "0.003 px"],
            ["Stage 3 (primitive fitting)", "1.32 px", "0.83 px"],
            ["→ all-views official d2c_eval (paired, 800 samples)", "1.448 px mean", "0.884 px mean (−39%)"],
        ],
        col_widths=[5.6, 3.8, 3.8], font_size=15, highlight_col=2,
    )

    # ---- Hatching deep dive ---------------------------------------------------
    section_slide(prs, "02", "Hatching — Deep Dive", "The hardest remaining problem, and what's in progress right now")

    bullet_slide(prs, "Why hatching is still the focus", "~45% of all primitives in patent drawings", [
        ("The side-layer fix above keeps hatches from breaking long outlines, but the "
         "*detector that decides hatch vs. not-hatch* still has poor precision.", 0),
        ("Concretely: it flagged 581 false \"hatch\" lines on a drill-view drawing that has "
         "no section hatching at all — and separately misses real diagonal hatching (poor recall).", 0),
        ("Five hand-crafted discriminators were tested to separate genuine hatching from "
         "dense mechanical line-work — spacing-CV, line density, FFT periodicity, gradient "
         "orientation-coherence, 2D fill-uniformity — and none separate the classes "
         "individually.", 0),
        ("Conclusion: heuristics can't fix this on their own → a learned detector is needed "
         "before the parametric HATCH-region rendering can be safely re-enabled.", 0, ),
    ], size=18)

    bullet_slide(prs, "The plan: two phases, cheapest first", None, [
        ("Phase 1 — region re-classifier (in progress)", 0),
        ("Keep the existing detector as a proposal generator; train a classifier on a joint "
         "feature vector (orientation entropy, FFT power, fill-uniformity, spacing-CV, "
         "enclosure) to accept/reject each candidate region.", 1),
        ("De-risk gate: by-figure-split ROC-AUC ≥ ~0.90 on labeled pilot data ⇒ viable.", 1),
        ("Phase 2 — pixel/patch segmentation CNN (fallback, only if Phase 1's recall is insufficient)", 0),
        ("A small U-Net/patch classifier on the cleaned image, independent of the proposal "
         "detector — fixes precision and recall together, but needs pixel-mask labels and "
         "GPU training (heavier).", 1),
    ], size=18)

    bullet_slide(prs, "Labeling tooling built this sprint", "tools/hatch_*.py", [
        ("hatch_features.py / hatch_label_sample.py — per-region feature vectors + a sampled "
         "pilot pack (34 figures, 213 candidate regions, overlays rendered).", 0),
        ("hatch_label.py — interactive labeler for the de-risk gate: one zoomed crop per "
         "candidate region, single y/n keypress (no IDs, no feature numbers to interpret).", 0),
        ("hatch_mask_label.py — independent ground-truth tool: draw the *actual* hatch areas "
         "directly (not just judging the detector's guesses), so it also catches hatching the "
         "detector misses entirely.", 0),
        ("Polygon mode (edge-snapped vertices) + circle mode (for tube/ring cross-sections), "
         "auto-split multi-panel sheets, scroll-zoom + drag-pan, candidate-region guide overlay.", 1),
        ("hatch_review.py — read-only QA viewer: detector proposals (cyan) vs. labeled ground "
         "truth (green) side by side, to audit a student's labeling pass before training.", 0),
        ("hatch_train.py — by-figure cross-validated ROC-AUC de-risk gate (HistGradientBoosting).", 0),
    ], size=16.5)

    bullet_slide(prs, "Status right now", None, [
        ("Ground-truth pass: 70 sub-drawings labeled, 40 hatch areas drawn (one student pass).", 0),
        ("Review pass just completed — some mislabeling found, to be corrected before training.", 0),
        ("Not yet done: fix the flagged mislabels, run the de-risk AUC, decide Phase 1 vs. "
         "Phase 2, integrate the winning classifier into Stage 2, and replace the convex-hull "
         "boundary with a concave hull / contour clip so HATCH regions don't over-cover "
         "concave or disjoint fills.", 0),
    ], size=19)

    # ---- Next steps -----------------------------------------------------------
    section_slide(prs, "03", "Known Limitations & Next Steps", "What's open across the whole pipeline")

    bullet_slide(prs, "Open items — ranked by what's next", None, [
        ("1. Hatching classifier (above) — fix labels, de-risk AUC, integrate, re-enable HATCH regions.", 0),
        ("2. Patent corner gap — fusion keypoints win on Drawing2CAD but are neutral on real patents "
         "(out-of-distribution). Needs a gold corner-labelled patent eval set (~50–100 sketches) and "
         "real-patent pseudo-labels for fine-tuning.", 0),
        ("3. Visual curation layer — scale from the current small CLIP-vetted seed (259 examples) to an "
         "active-learning classifier that recovers safe mechanical drawings the strict hand rules reject.", 0),
        ("4. Thin outline rectangles / concentric rings — skeletonize into two parallel \"rails\"; naive "
         "junction merging would weld them together. Needs a ladder-aware rung-removal.", 0),
        ("5. Stage 0 refinement — separate dashed/hidden construction lines from annotation/dimension "
         "geometry more reliably; add OCR so DXF reinjects real MTEXT, not just crop overlays.", 0),
        ("6. Runtime tail — dense hatch/detail sketches can still exceed the time budget; a per-sketch "
         "max_edges guard is still worth adding for full-corpus runs.", 0),
        ("7. Thin-stroke 1px position bias — Zhang-Suen skeletonization is systematically off-center for "
         "very thin strokes (radius ≤3px), a ~15–20% IoU hit on those shapes specifically.", 0),
    ], size=16)

    closing_slide(prs)

    _set_speaker_notes(prs)

    prs.save(out_path)
    return len(prs.slides._sldIdLst)


# One entry per slide, in creation order — see NOTES_ORDER comment above each
# block in build() if you reorder/add slides; keep this list in sync.
_NOTES = [
    # 1. Title
    "Open with the headline framing, not the title text: this pipeline turns a hand-drawn or "
    "patent engineering sketch into an editable SVG/DXF, and today is a status update on AP3 — "
    "what's been fixed, what the numbers say, and what's still open. Set expectations: most of "
    "the talk is problem/fix pairs with real before/after evidence, not just a feature list.",
    # 2. Agenda
    "Walk the agenda quickly — don't read it verbatim. Flag up front that hatching gets its own "
    "deep-dive at the end because it's the single largest unsolved problem (~45% of all "
    "primitives in patent drawings) and it's where the team's effort is concentrated right now.",
    # 3. Pipeline diagram
    "Anchor everyone on the five stages before diving into fixes, since every fix that follows "
    "will be labeled by which stage it lives in. Mention the two callouts: Stage 0 and Stage 2 "
    "both reinject material later (references, hachures) rather than discarding it — that "
    "reinjection design is why later slides can say 'excluded from confidence gates' without it "
    "meaning 'thrown away'. The 59-commits point is just to signal this is a maturing pipeline, "
    "not a one-shot prototype.",
    # 4. Section 01
    "Transition slide — give it one beat. Tell the room the structure: problem, fix, measured "
    "result, for six categories. Say explicitly that every number on the next slides comes from "
    "an actual eval run (Drawing2CAD and/or PatentData), not an estimate.",
    # 5. Corners
    "This is the keypoint-CNN work. The key nuance to state out loud: retraining a detector "
    "changes nothing by itself — topology extraction had to be rewired to actually consume the "
    "keypoints, which is why this was real engineering, not just a training run. Land the "
    "headline (fusion beats pure-CN on every D2C metric) but be upfront about the open gap on "
    "real patents — it's flagged again in Next Steps, so don't oversell it here.",
    # 6. Short strokes
    "This is the _simplify_graph fix — pure graph cleanup, no geometry moved, which is why "
    "accuracy is flat-to-better rather than a tradeoff. The mental model worth stating: the "
    "classical skeleton tracer creates a phantom junction at every staircase bend, so a single "
    "straight line becomes dozens of tiny edges; this pass merges them back. The next slide "
    "shows this on a real run.",
    # 7. Short strokes image
    "This is a real run on the repo's own sample image, not a mockup — same input skeleton both "
    "times, only the simplify_graph flag changed. Point at the screw thread and the big circle: "
    "before, they're confetti; after, they're each one or two consolidated strokes. 936 to 311 "
    "edges on this one sample.",
    # 8. Long strokes
    "Different root cause from the previous slide, same symptom (fragmentation) — worth saying "
    "explicitly so it doesn't sound redundant. This one is patent-specific: reference numerals "
    "and leader lines physically touch the real outlines, so removing them before Stage 1 even "
    "sees the image prevents the contact-point junctions from ever being created. Mention the "
    "OCR pivot fixed a real production incident (65% over-removal) before it became the default.",
    # 9. Curves
    "Two separate curve bugs fixed together, worth distinguishing: (1) splines wiggling off "
    "long straight lines — a guard rejects and falls back to RDP points; (2) compound curves "
    "(fillets, isometric silhouettes) that no single arc/ellipse could fit — now corner-split "
    "and fit piecewise. The 39% Chamfer drop is the single biggest accuracy jump in the whole "
    "project from one ~10-line fix, so let that number land.",
    # 10. Circles & polygons
    "Tell this as a two-steps-forward-one-step-back story: fixing spurious noise circles "
    "initially broke genuine small circles, because the first fix was a blunt size cutoff. The "
    "real fix is shape-aware (circularity, not size). Same lesson for rectangles — circle/ellipse "
    "RANSAC can't fit a right angle, so there's now a dedicated polygon fitter tried first.",
    # 11. Hatching slide (problem/fix card)
    "Set up the deep-dive: this card is the side-layer *topology* fix (keeps hatches from "
    "breaking outlines), which is different from and a prerequisite to the *detection* problem "
    "covered later. Be precise about that distinction or the audience will think hatching is "
    "fully solved. The next slide shows this fix on a real patent figure.",
    # 12. Hatching image
    "Real patent figure, not synthetic. Walk the four panels left to right: input, the "
    "fragmented old graph (141 edges, hatching cutting through outlines), the extracted hachure "
    "cluster on its own (113 edges, shown in magenta), and the clean result (25 edges). This is "
    "the side-layer mechanism made visible.",
    # 13. Headline metrics
    "These are the numbers that matter if someone only remembers one slide: 75% Chamfer "
    "reduction, zero-output samples eliminated entirely, and a 23% speedup despite all the new "
    "correctness checks. Note these are paired before/after on the identical 1000-sample test "
    "set, seed-locked — not cherry-picked.",
    # 14. Error attribution
    "This table is *why* the team knew where to spend effort: Stage 1 and Stage 2 contribute "
    "essentially zero geometric error — it's almost entirely Stage 3 fitting. That finding is "
    "what justified the curve-fitting investment on the earlier slide rather than, say, more "
    "preprocessing work. The isometric view's outsized drop (2.35 to 0.82) is the single "
    "clearest proof the B-spline fix mattered.",
    # 15. Section 02
    "Transition into the hatching deep-dive. Tone shift here: everything before this was "
    "*done*; this section is *in progress*, started this sprint, and the labeling tooling is "
    "what most of the recent work has actually been.",
    # 16. Why hatching
    "The 581-false-positive number is the concrete, memorable proof that heuristics failed — "
    "use it if anyone asks 'why not just tune the existing detector better.' The five-discriminator "
    "result is the answer to 'did you try simpler things first': yes, individually none of them "
    "separate the classes, which is exactly why a learned joint classifier is the next move.",
    # 17. The plan
    "Frame this as a deliberately cheap-first decision, not indecision: Phase 1 reuses the "
    "existing detector as a proposal generator and only adds a classifier on top, which is far "
    "less labeling than Phase 2's pixel masks. The AUC threshold (0.90) is a pre-committed "
    "go/no-go gate, decided before doing the labeling — so the next phase isn't a judgment call "
    "made after the fact.",
    # 18. Labeling tooling
    "This is mostly infrastructure built this sprint, worth walking through briefly tool by "
    "tool since it explains where the time went: feature extraction and sampling, the simplified "
    "per-region yes/no labeler, the from-scratch ground-truth drawing tool (polygon + circle "
    "modes, this is the one with the most iteration — zoom, pan, edge-snapping, multi-panel "
    "splitting all got built in response to real usability problems hit while labeling), and "
    "the QA review tool that's the reason mislabels were caught before training instead of after.",
    # 19. Status right now
    "Be direct about where this actually stands: labeling is done for this pass, QA just ran and "
    "found real mistakes, and nothing downstream (AUC, classifier, integration) has started yet. "
    "This is the honest 'what's blocking what happens next' slide — don't let it sound further "
    "along than it is.",
    # 20. Section 03
    "Transition to the full cross-pipeline limitations list — this is everything open, not just "
    "hatching, prioritized by what's actually next rather than by category.",
    # 21. Open items
    "These are ranked, not just listed — say that explicitly. Item 1 and 2 are the live "
    "priorities; items 3 through 7 are real but lower urgency. If pressed for a single 'what's "
    "next,' it's fixing the mislabels found in QA, then the de-risk AUC run.",
    # 22. Closing
    "Invite questions. If asked for sources, every number in this deck traces back to README.md's "
    "bug-fix log and project-status section, or to this sprint's hatching work — nothing here is "
    "an estimate.",
]


def _set_speaker_notes(prs):
    for slide, note in zip(prs.slides, _NOTES):
        slide.notes_slide.notes_text_frame.text = note


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--out", default="output/presentations/AP3_status.pptx")
    ap.add_argument("--figs", default="output/presentations/figs")
    a = ap.parse_args()
    import os
    os.makedirs(os.path.dirname(a.out), exist_ok=True)
    from pathlib import Path
    n = build(a.out, Path(a.figs))
    print(f"wrote {n} slides -> {a.out}")


if __name__ == "__main__":
    main()
